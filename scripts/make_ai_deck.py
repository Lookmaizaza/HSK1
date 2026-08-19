"""Build an AI-architecture pptx for the 語 ปากจีน app."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ---------------------------------------------------------------
BG = RGBColor(0xF7, 0xF7, 0xFA)
INK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x7C, 0x3A, 0xED)         # violet — Pathumma
ACCENT_SOFT = RGBColor(0xE9, 0xD5, 0xFF)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xE6, 0xE6, 0xFA)

# Use a slightly-wider 16:9 canvas.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]


# ---- helpers ---------------------------------------------------------------
def add_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # send to back — index 0
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Helvetica"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_pill(slide, x, y, w, h, text, *, fill=ACCENT, fg=RGBColor(0xFF, 0xFF, 0xFF), size=11):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.5
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(40000)
    tf.margin_top = tf.margin_bottom = Emu(20000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return box


def add_card(slide, x, y, w, h, *, fill=RGBColor(0xFF, 0xFF, 0xFF), border=RGBColor(0xE5, 0xE7, 0xEB)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    return card


def add_code(slide, x, y, w, h, code_text, *, size=12):
    add_card(slide, x, y, w, h, fill=CODE_BG, border=CODE_BG)
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15),
                                  w - Inches(0.4), h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(code_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = "Menlo"
        run.font.size = Pt(size)
        run.font.color.rgb = CODE_FG


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def header(slide, eyebrow, title, *, eyebrow_color=ACCENT):
    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.3),
             eyebrow.upper(), size=11, bold=True, color=eyebrow_color)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=INK)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.6), Inches(1.45),
                                        Inches(0.7), Inches(0.06))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()


def footer(slide, page_n, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "語 ปากจีน · AI architecture",
             size=10, color=MUTED)
    add_text(slide, Inches(11.6), Inches(7.05), Inches(1.2), Inches(0.3),
             f"{page_n} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ---- slide builders --------------------------------------------------------
def slide_title():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s, RGBColor(0x12, 0x06, 0x29))
    # gradient-ish accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               0, Inches(3.4), SLIDE_W, Inches(0.08))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    add_text(s, Inches(0.6), Inches(0.5), Inches(2), Inches(0.5),
             "AI ARCHITECTURE", size=12, bold=True, color=ACCENT_SOFT)
    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.5),
             "How the AI actually works", size=54, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, Inches(0.6), Inches(3.7), Inches(12), Inches(1.0),
             "語 ปากจีน — Mandarin pronunciation coach",
             size=22, color=ACCENT_SOFT)
    add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.6),
             "Powered by Pathumma ThaiLLM (NECTEC)",
             size=16, color=RGBColor(0xC4, 0xB5, 0xFD))


def slide_big_idea():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "the trick", "The AI never hears your voice")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(1.2),
             "Speech-to-text runs in your browser via the Web Speech API.\n"
             "By the time the request reaches the server, audio is gone — only the\n"
             "recognized hanzi text travels to the LLM.",
             size=18, color=INK)

    # Pipeline visual
    box_y = Inches(4.0)
    box_h = Inches(1.0)
    cols = [
        ("🎤 Voice", "(audio)", ROSE),
        ("Web Speech API", "zh-CN, on-device", AMBER),
        ('"你好"', "text only", EMERALD),
        ("Pathumma", "LLM judge", ACCENT),
        ("✓ feedback", "Thai", INK),
    ]
    col_w = Inches(2.3)
    gap = Inches(0.15)
    start_x = Inches(0.45)
    for i, (title, sub, color) in enumerate(cols):
        x = start_x + (col_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_y, col_w, box_h)
        card.adjustments[0] = 0.12
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_text(s, x, box_y + Inches(0.18), col_w, Inches(0.5),
                 title, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, x, box_y + Inches(0.6), col_w, Inches(0.4),
                 sub, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(cols) - 1:
            arrow_x1 = x + col_w
            arrow_x2 = arrow_x1 + gap
            arrow_y = box_y + box_h / 2
            add_arrow(s, arrow_x1, arrow_y, arrow_x2, arrow_y, color=MUTED, width=2)

    add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.0),
             "→ Pathumma is a judge that compares two strings (target vs heard),\n"
             "    not an audio model. Cheaper, faster, and works surprisingly well.",
             size=14, color=MUTED)


def slide_request_shape():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "step 1 · request", "What the lesson page sends")

    add_text(s, Inches(0.6), Inches(1.9), Inches(7), Inches(0.5),
             "POST /api/evaluate", size=16, bold=True, color=INK)
    add_code(s, Inches(0.6), Inches(2.4), Inches(7.5), Inches(2.4), """{
  "target": {
    "hanzi":   "你好",
    "pinyin":  "nǐ hǎo",
    "english": "Hello"
  },
  "said":   "你好",
  "tone":   3,
  "model":  "pathumma-thaillm"
}""", size=14)

    # Right-side annotations
    bx = Inches(8.5)
    add_text(s, bx, Inches(2.4), Inches(4.5), Inches(0.4),
             "target", size=14, bold=True, color=ACCENT)
    add_text(s, bx, Inches(2.7), Inches(4.5), Inches(0.6),
             "what you were asked to say",
             size=12, color=MUTED)

    add_text(s, bx, Inches(3.35), Inches(4.5), Inches(0.4),
             "said", size=14, bold=True, color=ACCENT)
    add_text(s, bx, Inches(3.65), Inches(4.5), Inches(0.6),
             "what the recognizer thinks you said",
             size=12, color=MUTED)

    add_text(s, bx, Inches(4.3), Inches(4.5), Inches(0.4),
             "tone", size=14, bold=True, color=ACCENT)
    add_text(s, bx, Inches(4.6), Inches(4.5), Inches(0.6),
             "only set for tone-drill lessons (2 or 3)",
             size=12, color=MUTED)

    add_text(s, bx, Inches(5.25), Inches(4.5), Inches(0.4),
             "model", size=14, bold=True, color=ACCENT)
    add_text(s, bx, Inches(5.55), Inches(4.5), Inches(0.6),
             "default = pathumma-thaillm",
             size=12, color=MUTED)


def slide_system_prompt():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "step 2 · prompt engineering",
           "Three behaviors hard-wired into the system prompt")

    # Three columns
    cols = [
        ("FORGIVING",
         "Tolerate recognizer slip-ups:\nmissed particles, near-homophones,\ndropped neutral tones.",
         EMERALD),
        ("STRICT JSON",
         "Reply ONLY with a one-line\nJSON object: { correct, score,\nfeedback, corrected }.",
         AMBER),
        ("THAI ONLY",
         "Feedback must be in Thai.\nShort, encouraging, points to\nwhat to fix.",
         ACCENT),
    ]
    col_w = Inches(3.9)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    top = Inches(2.0)
    for i, (label, body, color) in enumerate(cols):
        x = start_x + (col_w + gap) * i
        add_card(s, x, top, col_w, Inches(2.8))
        add_pill(s, x + Inches(0.3), top + Inches(0.3), Inches(1.6), Inches(0.4),
                 label, fill=color, size=11)
        add_text(s, x + Inches(0.3), top + Inches(1.0), col_w - Inches(0.6), Inches(1.8),
                 body, size=14, color=INK)

    add_text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
             "+ optional tone instruction (next slide) →",
             size=14, color=MUTED, bold=True)


def slide_tone_mode():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "step 3 · tone mode",
           "One extra block flips grading into tone-strict")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.6),
             "When the lesson has tone: 2 or tone: 3, toneInstruction() appends:",
             size=15, color=INK)

    add_code(s, Inches(0.6), Inches(2.5), Inches(12), Inches(1.8),
             "โหมดฝึกวรรณยุกต์: ผู้เรียนกำลังฝึก วรรณยุกต์ที่ 3 (เสียงตก-ขึ้น)\n"
             "- ตรวจวรรณยุกต์อย่างเข้มงวด\n"
             "- ถ้า hanzi ถูกแต่วรรณยุกต์ผิด → mark incorrect\n"
             "- จำกฎ tone sandhi: 3+3 → 2+3", size=14)

    add_card(s, Inches(0.6), Inches(4.7), Inches(12), Inches(1.9),
             fill=ACCENT_SOFT, border=ACCENT)
    add_text(s, Inches(0.85), Inches(4.85), Inches(3), Inches(0.4),
             "TRICK", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.85), Inches(5.2), Inches(11.5), Inches(1.5),
             "If you say the WRONG tone, the recognizer often returns a different\n"
             "hanzi that sounds like what you actually said.\n"
             "Example: target = 好 (hǎo). You say háo by mistake → recognizer hears 号.\n"
             "The LLM sees target=好, said=号, same syllable / different tone → catches it.",
             size=14, color=INK)


def slide_call_pathumma():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "step 4 · LLM call", "Posting to Pathumma")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
             "POST https://thaillm.or.th/api/v1/chat/completions",
             size=14, bold=True, color=INK)
    add_text(s, Inches(0.6), Inches(2.35), Inches(12), Inches(0.5),
             "Authorization: Bearer THAILLM_API_KEY",
             size=13, color=MUTED)

    add_code(s, Inches(0.6), Inches(3.0), Inches(12), Inches(3.6), """{
  "model": "pathumma-thaillm-qwen3-8b-think-3.0.0",
  "messages": [
    { "role": "system", "content": "<base prompt> + <tone instruction>" },
    { "role": "user",   "content": "TARGET hanzi: 你好\\n
                                    TARGET pinyin: nǐ hǎo\\n
                                    TARGET english: Hello\\n
                                    HEARD: 你好" }
  ],
  "temperature": 0.3,
  "max_tokens": 2048
}""", size=13)


def slide_think_problem():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "step 5 · reasoning model",
           "Why Pathumma's reply needs extra parsing")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.6),
             "Pathumma is a *reasoning* model (qwen3-8b-think).\n"
             "It emits chain-of-thought wrapped in <think>…</think> tags before the answer:",
             size=14, color=INK)

    add_code(s, Inches(0.6), Inches(3.0), Inches(12), Inches(2.3),
             "<think>\n"
             "ผู้ใช้พูด 你好 ตรงกับ target พอดี วรรณยุกต์ 3+3 → 2+3 ถูกต้อง\n"
             "ให้คะแนน 100 ได้เลย\n"
             "</think>\n"
             '{ "correct": true, "score": 100, "feedback": "ยอดเยี่ยม!", "corrected": null }',
             size=13)

    add_card(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.1),
             fill=RGBColor(0xFE, 0xE2, 0xE2), border=ROSE)
    add_text(s, Inches(0.85), Inches(5.75), Inches(11.5), Inches(0.9),
             "JSON.parse on this raw string → crashes. We need to strip <think>…</think> first.",
             size=14, color=INK, bold=True)


def slide_json_extract():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "step 6 · json extraction",
           "extractJsonObject() — three-stage cleanup")

    steps = [
        ("1", "Strip <think>…</think> blocks",
         "Regex: /<think>[\\s\\S]*?<\\/think>/gi"),
        ("2", "Strip markdown code fences",
         "Regex: /```(?:json)?\\s*([\\s\\S]*?)```/i"),
        ("3", "Find first balanced { … } block",
         "Walk the string, count braces, stop when depth == 0"),
    ]
    y = Inches(2.0)
    for i, (n, title, sub) in enumerate(steps):
        cy = y + Inches(i * 1.5)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.7), cy, Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        add_text(s, Inches(0.7), cy + Inches(0.15), Inches(0.8), Inches(0.5),
                 n, size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.8), cy, Inches(11), Inches(0.5),
                 title, size=18, bold=True, color=INK)
        add_text(s, Inches(1.8), cy + Inches(0.5), Inches(11), Inches(0.5),
                 sub, size=13, color=MUTED, font="Menlo")

    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             "Robust whether Pathumma emits <think> tags or Gemini emits pure JSON.",
             size=13, color=MUTED)


def slide_fallback():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "step 7 · fallback",
           "When the LLM has a bad day")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.6),
             "If extraction or parsing fails, we don't crash — we degrade gracefully:",
             size=15, color=INK)

    add_code(s, Inches(0.6), Inches(2.6), Inches(12), Inches(2.5),
             "const anyMatch = [...target.hanzi].some(c => said.includes(c));\n"
             "parsed = {\n"
             "  correct:   anyMatch,\n"
             "  score:     anyMatch ? 60 : 0,\n"
             '  feedback:  anyMatch ? "ใกล้แล้ว!" : "ลองอีกครั้งนะ",\n'
             "  corrected: anyMatch ? null : target.hanzi\n"
             "};", size=13)

    add_card(s, Inches(0.6), Inches(5.4), Inches(12), Inches(1.5),
             fill=RGBColor(0xDC, 0xFC, 0xE7), border=EMERALD)
    add_text(s, Inches(0.85), Inches(5.55), Inches(11.5), Inches(1.3),
             "Dumb character-overlap heuristic = lenient but keeps the UI flowing.\n"
             "Better to give a forgiving 60/100 than break the lesson with a 500.",
             size=14, color=INK)


def slide_response():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "step 8 · response",
           "What comes back, what the UI does")

    add_code(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(2.5), """{
  "correct":   true,
  "score":     100,
  "feedback":  "ยอดเยี่ยมมาก!",
  "corrected": null
}""", size=14)

    bx = Inches(7.4)
    add_card(s, bx, Inches(2.0), Inches(5.4), Inches(1.1),
             fill=RGBColor(0xDC, 0xFC, 0xE7), border=EMERALD)
    add_text(s, bx + Inches(0.3), Inches(2.15), Inches(5), Inches(0.4),
             "correct = true", size=14, bold=True, color=EMERALD)
    add_text(s, bx + Inches(0.3), Inches(2.5), Inches(5), Inches(0.6),
             "→ green card · +10 XP · advance",
             size=13, color=INK)

    add_card(s, bx, Inches(3.3), Inches(5.4), Inches(1.1),
             fill=RGBColor(0xFE, 0xE2, 0xE2), border=ROSE)
    add_text(s, bx + Inches(0.3), Inches(3.45), Inches(5), Inches(0.4),
             "correct = false", size=14, bold=True, color=ROSE)
    add_text(s, bx + Inches(0.3), Inches(3.8), Inches(5), Inches(0.6),
             "→ red card · -1 heart · show corrected hanzi",
             size=13, color=INK)

    add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.5),
             "If user is logged in, the XP/heart change is also posted to Turso via /api/progress.",
             size=13, color=MUTED)


def slide_why():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    header(s, "design rationale", "Why LLM-as-judge instead of…")

    rows = [
        ("Whisper / audio LLM",
         "Hears real tones — but 10× cost, 2-3s latency.",
         ROSE),
        ("Plain string equality",
         "Marks you wrong on every near-homophone. Brutal UX.",
         ROSE),
        ("Pure rule-based grader",
         "Can't give natural Thai feedback or explain why.",
         ROSE),
        ("LLM-as-judge (ours)",
         "Fast · cheap · forgiving where it should be · strict when tone mode is on · Thai feedback for free.",
         EMERALD),
    ]
    y = Inches(2.0)
    for i, (title, body, color) in enumerate(rows):
        cy = y + Inches(i * 1.15)
        add_card(s, Inches(0.6), cy, Inches(12.0), Inches(0.95),
                 border=color)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.85), cy + Inches(0.35),
                                 Inches(0.25), Inches(0.25))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_text(s, Inches(1.3), cy + Inches(0.15), Inches(11), Inches(0.4),
                 title, size=16, bold=True, color=INK)
        add_text(s, Inches(1.3), cy + Inches(0.5), Inches(11), Inches(0.4),
                 body, size=13, color=MUTED)


def slide_end():
    s = prs.slides.add_slide(blank_layout)
    add_bg(s, RGBColor(0x12, 0x06, 0x29))
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.5),
             "AI ≠ audio model.",
             size=56, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, Inches(0.6), Inches(3.6), Inches(12), Inches(1.5),
             "AI = text judge with a strict prompt.",
             size=36, bold=True, color=ACCENT_SOFT)
    add_text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.8),
             "Browser does the listening · Pathumma does the grading · Thai feedback for free.",
             size=18, color=RGBColor(0xC4, 0xB5, 0xFD))


# ---- build ----------------------------------------------------------------
builders = [
    slide_title,
    slide_big_idea,
    slide_request_shape,
    slide_system_prompt,
    slide_tone_mode,
    slide_call_pathumma,
    slide_think_problem,
    slide_json_extract,
    slide_fallback,
    slide_response,
    slide_why,
    slide_end,
]
total = len(builders)
for fn in builders:
    fn()

# add footers to all but title + end
for i in range(1, total - 1):
    footer(prs.slides[i], i + 1, total)

out_path = "ai_architecture.pptx"
prs.save(out_path)
print(f"wrote {out_path} — {total} slides")
